from __future__ import annotations

from pathlib import Path

from docx import Document as WordDocument
from pptx import Presentation
from pypdf import PdfReader

from harness_app.knowledge.models import ExtractionResult
from harness_app.knowledge.ocr import OcrFallback


class DocumentTextExtractor:
    def __init__(self, ocr_fallback: OcrFallback) -> None:
        self._ocr_fallback = ocr_fallback

    def extract(self, path: Path) -> ExtractionResult:
        extension = path.suffix.lower()
        if extension == ".docx":
            return self._extract_docx(path)
        if extension == ".pptx":
            return self._extract_pptx(path)
        if extension == ".pdf":
            return self._extract_pdf(path)
        return ExtractionResult(content="", notes=[f"Unsupported file extension: {extension}"], used_ocr=False)

    def _extract_docx(self, path: Path) -> ExtractionResult:
        document = WordDocument(path)
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        content = "\n".join(paragraphs).strip()
        if content:
            return ExtractionResult(content=content, notes=["Extracted text from DOCX paragraphs."])
        fallback = self._ocr_fallback.extract_text(path)
        fallback.notes.insert(0, "DOCX text extraction returned no content.")
        return fallback

    def _extract_pptx(self, path: Path) -> ExtractionResult:
        presentation = Presentation(path)
        text_runs: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    text_runs.append(f"[Slide {slide_index}] {text.strip()}")
        content = "\n".join(text_runs).strip()
        if content:
            return ExtractionResult(content=content, notes=["Extracted text from PPTX slide shapes."])
        fallback = self._ocr_fallback.extract_text(path)
        fallback.notes.insert(0, "PPTX text extraction returned no content.")
        return fallback

    def _extract_pdf(self, path: Path) -> ExtractionResult:
        reader = PdfReader(str(path))
        page_text: list[str] = []
        for page_index, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                page_text.append(f"[Page {page_index}] {text}")
        content = "\n".join(page_text).strip()
        if content:
            return ExtractionResult(content=content, notes=["Extracted selectable text from PDF pages."])
        fallback = self._ocr_fallback.extract_text(path)
        fallback.notes.insert(0, "PDF text extraction returned no content.")
        return fallback
